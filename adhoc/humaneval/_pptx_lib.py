"""
Shared PPTX helpers for human-eval PPTX scripts under dev_robosuite/adhoc/humaneval.
Workspace = repo root: parents[3] from dev_robosuite/adhoc/humaneval/<script>.py
"""
from __future__ import annotations

import json
import math
import os
import platform
import random
import shlex
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Literal

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# dev_robosuite/adhoc/humaneval/foo.py -> parents[3]=workspace
WS = Path(__file__).resolve().parents[3]

# Vivid on-screen / PPT (chips, A–D panel keys when mark_answer)
ANSWER_KEY_FONT_RGB: tuple[int, int, int] = (0, 145, 255)
# compare_baseline: sophisticated (reference policy) when marking
SOPH_KEY_GREEN_RGB: tuple[int, int, int] = (0, 168, 72)
DEV_R = WS / "dev_robosuite"
MOTIONS = DEV_R / "data" / "motions"
SEED = DEV_R / "data" / "seed"
_MC_MANIP = DEV_R / "data" / "results" / "motion_configs" / "manipulator"
_MC_BUNDLES = _MC_MANIP / "bundles"
_HUMAN_ROOT = DEV_R / "data" / "results" / "human_eval"

HUMAN = {
    "manipulator": {
        "sophisticated_dir": DEV_R / "data" / "human_eval" / "sophisticated_v1",
    },
    "tiago": {
        "catalog": _HUMAN_ROOT / "mobile_mani_v1" / "catalog.json",
        "binary": _HUMAN_ROOT / "mobile_mani_v1" / "binary_items.json",
    },
    "quadruped": {
        "catalog": _HUMAN_ROOT / "quadruped_v1" / "catalog.json",
        "binary": _HUMAN_ROOT / "quadruped_v1" / "binary_items.json",
    },
}

MANIP_CFG = {
    "iconic": {
        "soph": _MC_MANIP / "motion_configs_prompt_v19_sophisticated.json",
        "nr": _MC_BUNDLES
        / "baseline_prompt19_full_no_reasoning"
        / "motion_configs_prompt_v19_sophisticated_no_reasoning_iconic.json",
        "soph_gif": MOTIONS / "v19_sophisticated" / "IIWA",
        "nr_gif": MOTIONS / "baseline_prompt19_full_no_reasoning" / "no_reasoning_iconic" / "IIWA",
    },
    "contextual": {
        "soph": _MC_MANIP / "motion_configs_prompt_v19_sophisticated_contextual.json",
        "nr": _MC_BUNDLES
        / "baseline_prompt19_full_no_reasoning"
        / "motion_configs_prompt_v19_sophisticated_no_reasoning_contextual.json",
        "soph_gif": MOTIONS / "v19_sophisticated_contextual_q4filled" / "IIWA",
        "nr_gif": MOTIONS / "baseline_prompt19_full_no_reasoning" / "no_reasoning_contextual" / "IIWA",
    },
}


def workspace() -> Path:
    return WS


def open_output_file(path: Path | str, *, do_open: bool = True) -> None:
    """
    On success, open the PPTX in the system default app (e.g. macOS: `open path`).
    Set do_open=False to skip.
    """
    if not do_open:
        return
    p = Path(path)
    if not p.is_file():
        return
    q = shlex.quote(str(p.resolve()))
    s = platform.system()
    if s == "Darwin":
        os.system(f"open {q}")
    elif s == "Windows":
        # os.startfile is more reliable than `start` with quoting
        os.startfile(str(p.resolve()))  # type: ignore[attr-defined]
    else:
        os.system(f"xdg-open {q}")


def add_do_open_arg(parser) -> None:
    """
    --do_open / --no-do_open (default: open). On Python 3.9+ uses BooleanOptionalAction;
    older: --no_open only.
    """
    import argparse
    import sys

    if sys.version_info >= (3, 9):
        parser.add_argument(
            "--do_open",
            action=argparse.BooleanOptionalAction,
            default=True,
            help="Open the generated PPTX in the default viewer when done (default: on).",
        )
    else:
        parser.set_defaults(no_open=False)
        parser.add_argument(
            "--no_open",
            action="store_true",
            help="Do not open the generated PPTX in the default viewer at the end.",
        )


def do_open_effective_from_args(args) -> bool:
    import sys

    if sys.version_info >= (3, 9):
        return bool(getattr(args, "do_open", True))
    return not bool(getattr(args, "no_open", False))


def add_mark_answer_arg(parser) -> None:
    """
    --mark_answer / --no-mark_answer (default: off). When on, PPTX can show review keys
    (e.g. binary GT yes/no) where supported.
    """
    import argparse
    import sys

    if sys.version_info >= (3, 9):
        parser.add_argument(
            "--mark_answer",
            action=argparse.BooleanOptionalAction,
            default=False,
            help="Show answer key on slides (e.g. binary ground-truth yes/no) when True.",
        )
    else:
        parser.set_defaults(mark_answer=False)
        parser.add_argument(
            "--mark_answer",
            action="store_true",
            help="Show answer key on slides (e.g. binary ground-truth yes/no).",
        )


def mark_answer_effective_from_args(args) -> bool:
    return bool(getattr(args, "mark_answer", False))


def default_out_path(kind: str, robot: str) -> Path:
    """data/results/pptx/YYYYMMDD_<kind>_<robot>.pptx"""
    d = DEV_R / "data" / "results" / "pptx"
    d.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d")
    return d / f"{stamp}_{kind}_{robot}.pptx"


def parse_robot(robot: str) -> list[str]:
    s = (robot or "").strip().lower()
    order = ("manipulator", "tiago", "quadruped")
    if s == "all":
        return list(order)
    if "," in s:
        parts = [p.strip() for p in s.split(",") if p.strip()]
        for p in parts:
            if p not in order:
                raise ValueError(
                    f"robot must be one of: {', '.join(order)}, all, or a comma list (e.g. manipulator,tiago); got {p!r}"
                )
        return parts
    if s in order:
        return [s]
    raise ValueError(
        f"robot must be one of: {', '.join(order)}, all, or a comma list (e.g. manipulator,tiago); got {robot!r}"
    )


def load_json(p: Path) -> Any:
    return json.loads(p.read_text(encoding="utf-8"))


def safe_name(text: str) -> str:
    return str(text).replace("/", "_").replace("\\", "_").replace(" ", "_")


def latest_gif_in_dir(base: Path, cue: str) -> Path | None:
    safe = safe_name(cue)
    if not base.exists():
        return None
    matches = sorted(base.rglob(f"*_{safe}_p*.gif"), key=lambda p: p.stat().st_mtime)
    return matches[-1] if matches else None


def _latest_gif_with_cue_idx(base: Path, cue: str, idx: int) -> Path | None:
    """Same resolution order as run_prompt19_baseline_experiment._latest_gif."""
    safe = safe_name(cue)
    if not base.exists():
        return None
    tiled = sorted(base.rglob(f"*_{safe}_c{idx}_tiled.gif"))
    if tiled:
        return tiled[-1]
    single = sorted(base.rglob(f"*_{safe}_p*.gif"))
    if single:
        return single[-1]
    any_match = sorted(base.rglob(f"*_{safe}_*.gif"))
    return any_match[-1] if any_match else None


def latest_iiwa_direct_baseline_gif(cue: str, idx: int, subtest: str, *, joint: bool) -> Path | None:
    """
    Prompt-19 **direct few-shot** renders: joint keyframes vs Cartesian+xyz (direct_xyz_theta; RPY family).
    Scans baseline_prompt19_direct_experiment* trees (same roots as baseline HTML builders).
    """
    dataset = subtest if subtest in ("iconic", "contextual") else "iconic"
    robot = "IIWA"
    sub = "direct_joint" if joint else "direct_xyz_theta"
    roots = [
        MOTIONS / "baseline_prompt19_direct_experiment" / sub / dataset / robot,
        MOTIONS / "baseline_prompt19_direct_experiment_more20configs" / sub / dataset / robot,
        MOTIONS / "baseline_prompt19_direct_experiment_extra10" / sub / dataset / robot,
        MOTIONS / "baseline_prompt19_direct_experiment_nonoverlap10" / sub / dataset / robot,
    ]
    for base in roots:
        g = _latest_gif_with_cue_idx(base, cue, int(idx))
        if g and g.is_file():
            return g
    return None


def reanchor_stale_path_under_workspace(absolute: str) -> Path | None:
    """
    JSON may store /other/machine/.../dev_robosuite/data/... — rebuild as WS / tail from
    'dev_robosuite/' or 'dev_locomotion/'.
    """
    s = str(absolute).replace("\\", "/")
    for head in ("dev_robosuite/", "dev_locomotion/"):
        j = s.find(head)
        if j >= 0:
            q = WS / s[j:].lstrip("/")
            if q.is_file():
                return q
    return None


def resolve_binary_item_gif(item: dict, *, prefer_no_reasoning: bool = False) -> Path | None:
    """
    sophisticated_eval batch items: abs gif_path can point to a removed timestamp. Fall back
    to reanchor under WS, then latest *_{cue}_p*.gif by testset (q4filled for contextual).

    If ``prefer_no_reasoning`` is True (e.g. binary PPTX for IIWA), resolve under the
    no-reasoning render roots in ``MANIP_CFG``; no fallback to on-disk sophisticated paths
    in the JSON.
    """
    cue = (item.get("source_cue") or item.get("shown_cue") or "").strip()
    ts = (item.get("testset") or item.get("subtest") or "").lower()
    if prefer_no_reasoning and cue:
        if ts == "iconic":
            g = latest_gif_in_dir(MANIP_CFG["iconic"]["nr_gif"], cue)
            if g and g.is_file():
                return g
            return None
        if ts == "contextual":
            g = latest_gif_in_dir(MANIP_CFG["contextual"]["nr_gif"], cue)
            if g and g.is_file():
                return g
            return None
        return None
    p = item.get("gif_path")
    if p:
        q = Path(p)
        if q.is_file():
            return q
        ra = reanchor_stale_path_under_workspace(p)
        if ra:
            return ra
    rel = item.get("gif_workspace_rel")
    if rel:
        q2 = WS / rel
        if q2.is_file():
            return q2
    if not cue:
        return None
    if ts == "iconic":
        return latest_gif_in_dir(MOTIONS / "v19_sophisticated" / "IIWA", cue)
    if ts == "contextual":
        for root in (
            MOTIONS / "v19_sophisticated_contextual_q4filled" / "IIWA",
            MOTIONS / "v19_sophisticated_contextual" / "IIWA",
        ):
            g = latest_gif_in_dir(root, cue)
            if g:
                return g
    return None


# --- text / shapes (aligned with build_prompt19_goodset_labeling_pptx) ---


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 20,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_rgb: tuple[int, int, int] = (0, 0, 0),
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = "Arial"
    r.font.color.rgb = RGBColor(*font_rgb)
    return box


def add_choice_chip(slide, left, top, width, height, text: str, *, text_rgb: tuple[int, int, int] = (0, 0, 0)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    shape.line.color.rgb = RGBColor(180, 186, 196)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = RGBColor(*text_rgb)
    return shape


def fit_picture(slide, img_path: Path, left, top, width, height) -> bool:
    if not img_path.exists():
        return False
    with Image.open(img_path) as im:
        img_w, img_h = im.size
    box_ratio = float(width) / float(height)
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        draw_w = width
        draw_h = width / img_ratio
        draw_left = left
        draw_top = top + (height - draw_h) / 2
    else:
        draw_h = height
        draw_w = height * img_ratio
        draw_top = top
        draw_left = left + (width - draw_w) / 2
    slide.shapes.add_picture(str(img_path), draw_left, draw_top, width=draw_w, height=draw_h)
    return True


@dataclass
class ProductRow:
    subtest: str
    idx: int
    cue: str
    label: str
    gif: Path | None
    robot_key: str


def sort_key_subtest(s: str) -> int:
    s = (s or "").lower()
    if s == "iconic":
        return 0
    if s == "contextual":
        return 1
    return 2


ManipMotion = Literal["soph", "nr"]


def load_manipulator_product_list(*, motion: ManipMotion = "soph") -> list[ProductRow]:
    """
    Cue list still comes from the sophisticated config JSON; ``motion`` chooses which
    on-disk tree supplies ``ProductRow.gif`` (sophisticated vs. no-reasoning renders).
    """
    if motion not in ("soph", "nr"):
        raise ValueError("motion must be 'soph' or 'nr'")
    rows: list[ProductRow] = []
    for sub, spec in MANIP_CFG.items():
        for r in sorted(load_json(spec["soph"]), key=lambda x: int(x["idx"])):
            cue = r["cue"]
            gdir = spec["soph_gif"] if motion == "soph" else spec["nr_gif"]
            g = latest_gif_in_dir(gdir, cue)
            rows.append(
                ProductRow(
                    subtest=sub,
                    idx=int(r["idx"]),
                    cue=cue,
                    label=(r.get("description") or "")[:240],
                    gif=g,
                    robot_key="manipulator",
                )
            )
    return rows


def load_catalog_gifs(robot: str, *, manipulator_gif: ManipMotion = "soph") -> list[ProductRow]:
    if robot == "manipulator":
        return load_manipulator_product_list(motion=manipulator_gif)
    p = HUMAN[robot]["catalog"]
    data = load_json(p)
    rows: list[ProductRow] = []
    for it in data.get("items") or []:
        rel = it.get("gif_workspace_rel")
        g = (WS / rel) if rel else None
        if g and not g.exists():
            g = None
        rows.append(
            ProductRow(
                subtest=it.get("subtest", ""),
                idx=int(it.get("idx", 0)),
                cue=it.get("cue", ""),
                label=(it.get("cue_text") or "")[:240],
                gif=g,
                robot_key=robot,
            )
        )
    rows.sort(key=lambda x: (sort_key_subtest(x.subtest), x.idx, x.cue))
    return rows


def apply_sample_n(rows: list[ProductRow], sample_n: int | None) -> list[ProductRow]:
    if sample_n is None or len(rows) <= sample_n:
        return rows
    return rows[:sample_n]


def iiwa_row_has_all_four_baselines(r: ProductRow) -> bool:
    """True if soph + nr + joint few-shot + xyz few-shot GIFs all resolve (IIWA compare_baseline)."""
    if r.robot_key != "manipulator" or r.subtest not in MANIP_CFG:
        return False
    spec = MANIP_CFG[r.subtest]
    g_soph = r.gif
    if not g_soph or not g_soph.is_file():
        return False
    g_nr = latest_gif_in_dir(spec["nr_gif"], r.cue)
    if not g_nr or not g_nr.is_file():
        return False
    g_j = latest_iiwa_direct_baseline_gif(r.cue, r.idx, r.subtest, joint=True)
    g_x = latest_iiwa_direct_baseline_gif(r.cue, r.idx, r.subtest, joint=False)
    if not g_j or not g_j.is_file() or not g_x or not g_x.is_file():
        return False
    return True


def apply_sample_n_iiwa_fewshot_first(
    rows: list[ProductRow], sample_n: int | None, seed: int
) -> list[ProductRow]:
    """
    Prefer items whose direct few-shot (joint + xyz) and nr baselines all exist, then fill from the rest.
    Shuffles stably within each pool. When sample_n is None, returns all rows (unchanged order).
    """
    if sample_n is None or len(rows) <= sample_n:
        return rows
    rng = random.Random(int(seed) & 0x7FFFFFFF)
    full = [r for r in rows if iiwa_row_has_all_four_baselines(r)]
    rest = [r for r in rows if not iiwa_row_has_all_four_baselines(r)]
    rng.shuffle(full)
    rng.shuffle(rest)
    out = full[:sample_n]
    if len(out) < sample_n:
        out.extend(rest[: sample_n - len(out)])
    return out


def add_humaneval_appendix_slides(
    prs: Presentation,
    *,
    title: str,
    column_headers: list[str],
    data_rows: list[list[str]],
    total_slides_line: str | None = None,
    max_rows_per_slide: int = 14,
) -> int:
    """
    One or more blank slides with a table: run metadata for reviewers (slide #, cue, answer, etc.).
    Returns: number of slides added.
    """
    if not column_headers:
        return 0
    ncols = len(column_headers)
    n = len(data_rows)
    if n == 0:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(s, Inches(0.45), Inches(0.35), Inches(12.0), Inches(0.42), title, size=15, bold=True)
        if total_slides_line:
            add_textbox(s, Inches(0.45), Inches(0.85), Inches(12.0), Inches(0.35), total_slides_line, size=9)
        add_textbox(
            s, Inches(0.45), Inches(1.2), Inches(12.0), Inches(0.3),
            "(no item-level rows in this run — title/section slides only).",
            size=9,
        )
        return 1
    start = 0
    part = 0
    while start < n:
        chunk = data_rows[start : start + max_rows_per_slide]
        s = prs.slides.add_slide(prs.slide_layouts[6])
        y_top = 0.32
        add_textbox(
            s, Inches(0.45), Inches(y_top), Inches(12.0), Inches(0.42),
            title if part == 0 else f"{title} (cont.)", size=15, bold=True,
        )
        y_top += 0.45
        if total_slides_line and part == 0:
            add_textbox(
                s, Inches(0.45), Inches(y_top), Inches(12.0), Inches(0.28), total_slides_line, size=9,
            )
            y_top += 0.32
        if n > max_rows_per_slide:
            add_textbox(
                s, Inches(0.45), Inches(y_top), Inches(12.0), Inches(0.22),
                f"Rows {start + 1}–{start + len(chunk)} of {n}",
                size=8,
            )
            y_top += 0.26
        nrows = 1 + len(chunk)
        h_body = min(6.2, 0.4 + 0.22 * nrows)
        _tbl = s.shapes.add_table(
            nrows, ncols, Inches(0.4), Inches(y_top + 0.02), Inches(12.5), Inches(h_body)
        )
        table = _tbl.table
        for j, h in enumerate(column_headers):
            c = table.cell(0, j)
            c.text = h
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(8)
                    r.font.bold = True
        for i, row in enumerate(chunk):
            for j, val in enumerate(row):
                if j >= ncols:
                    break
                c = table.cell(1 + i, j)
                c.text = (val or "")[:500]
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(7)
        part += 1
        start += len(chunk)
    return part


def load_sophisticated_binary_items() -> list[dict]:
    d = HUMAN["manipulator"]["sophisticated_dir"]
    items: list[dict] = []
    for f in sorted(d.glob("sophisticated_eval_batch_*.json")):
        js = load_json(f)
        items.extend(js.get("items") or [])
    return items


def load_binary_for_robot(robot: str) -> list[dict]:
    if robot == "manipulator":
        return load_sophisticated_binary_items()
    return load_json(HUMAN[robot]["binary"])


def product_grid_layout(per_page: int) -> tuple[int, int]:
    """cols, rows for up to per_page cells."""
    p = max(1, int(per_page))
    cols = int(math.ceil(math.sqrt(p)))
    rows = int(math.ceil(p / cols))
    return cols, rows
