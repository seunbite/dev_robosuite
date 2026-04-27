#!/usr/bin/env python3
"""
Binary cue–motion classification PPTX (Shown cue + GIF + YES/NO chips).

Data:
  - manipulator: human_eval/sophisticated_v1/sophisticated_eval_batch_*.json (GIFs: no-reasoning IIWA trees)
  - tiago: mobile_mani_v1/binary_items.json
  - quadruped: quadruped_v1/binary_items.json

Output: data/results/pptx/YYYYMMDD_binary_classification_<robot>.pptx
"""
from __future__ import annotations

import argparse
import random
import sys
from pathlib import Path

_H = Path(__file__).resolve().parent
if str(_H) not in sys.path:
    sys.path.insert(0, str(_H))

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from _pptx_lib import (
    ANSWER_KEY_FONT_RGB,
    add_choice_chip,
    add_do_open_arg,
    add_humaneval_appendix_slides,
    add_mark_answer_arg,
    add_textbox,
    default_out_path,
    do_open_effective_from_args,
    fit_picture,
    load_binary_for_robot,
    mark_answer_effective_from_args,
    open_output_file,
    parse_robot,
    resolve_binary_item_gif,
)

_BLACK = RGBColor(0, 0, 0)

APPENDIX_HEADERS = [
    "Deck slide #",
    "Item #",
    "Robot / section",
    "Cue name",
    "Cue index",
    "Contextual",
    "Answer key",
]


def _header_line(robot_label: str, it: dict, sub: str, order: int, total: int) -> str:
    """No pair_id / paths: those look like filenames and spoil the task."""
    aid = (it.get("assignment_id") or "").strip()
    if aid:
        return f"{robot_label}  ·  {sub}  ·  {aid}"
    return f"{robot_label}  ·  {sub}  ·  item {order} / {total}"


def _build_title(prs, title: str, subtitle: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, Inches(0.6), Inches(0.45), Inches(12.0), Inches(0.55), title, size=22, bold=True)
    add_textbox(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(2.0), subtitle, size=14)


def _slide_item(
    prs,
    order: int,
    total: int,
    it: dict,
    robot_label: str,
    *,
    mark_answer: bool,
    appendix: list[list[str]],
    prefer_nr_gif: bool = False,
) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    sub = it.get("subtest") or it.get("testset", "")
    add_textbox(
        s, Inches(0.45), Inches(0.25), Inches(10.5), Inches(0.42),
        _header_line(robot_label, it, sub, order, total),
        size=11,
        bold=True,
    )
    add_textbox(
        s, Inches(0.55), Inches(0.72), Inches(12.0), Inches(0.4),
        "Does the shown cue match this motion?",
        size=15,
        bold=True,
    )
    shown = it.get("shown_cue", "")
    # Single line: "Shown cue: ..." only (no long natural-language description under it)
    box = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65), Inches(1.1), Inches(11.9), Inches(0.55),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    box.line.color.rgb = RGBColor(200, 206, 216)
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = f"Shown cue: {shown}"
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = _BLACK
    gif = resolve_binary_item_gif(it, prefer_no_reasoning=prefer_nr_gif)
    if gif:
        fit_picture(s, gif, Inches(2.0), Inches(1.9), Inches(9.3), Inches(3.75))
    else:
        add_textbox(s, Inches(2.0), Inches(2.5), Inches(9.0), Inches(0.4), "(GIF not found)", size=11)
    gt = (it.get("ground_truth") or "").strip().lower()
    if mark_answer and gt in ("yes", "no"):
        yes_rgb = ANSWER_KEY_FONT_RGB if gt == "yes" else (0, 0, 0)
        no_rgb = ANSWER_KEY_FONT_RGB if gt == "no" else (0, 0, 0)
    else:
        yes_rgb = no_rgb = (0, 0, 0)
    add_choice_chip(
        s, Inches(3.6), Inches(6.5), Inches(2.1), Inches(0.52), "YES", text_rgb=yes_rgb
    )
    add_choice_chip(
        s, Inches(7.1), Inches(6.5), Inches(2.1), Inches(0.52), "NO", text_rgb=no_rgb
    )
    add_textbox(s, Inches(11.2), Inches(6.65), Inches(1.4), Inches(0.25), f"{order}/{total}", size=9)
    cidx = str(it.get("config_idx", it.get("idx", it.get("cue_index", "—"))))
    ts = (it.get("testset") or it.get("subtest") or sub or "").lower()
    ctx = "yes" if ts == "contextual" else "no"
    ans = gt if gt in ("yes", "no") else (gt or "—")
    appendix.append(
        [
            str(len(prs.slides)),
            f"{order}/{total}",
            robot_label,
            (shown or it.get("source_cue") or it.get("cue") or "—")[:200],
            cidx,
            ctx,
            ans,
        ]
    )


def _labels(rk: str) -> str:
    return {
        "manipulator": "Manipulator (IIWA · no-reasoning · HEval v1)",
        "tiago": "TIAGo (mobile_mani_v1 · binary)",
        "quadruped": "Quadruped (locomotion v1 · binary)",
    }[rk]


def build(
    robot: str = "all",
    sample_n: int | None = None,
    seed: int = 17,
    out: str | None = None,
    mark_answer: bool = False,
) -> Path:
    rlist = parse_robot(robot)
    out_p = (
        Path(out)
        if out
        else default_out_path("binary_classification", "_".join(rlist) if len(rlist) > 1 else rlist[0])
    )
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sub = (
        "Task: Does the **shown** cue match the motion? YES / NO.\n"
        + (
            "With --mark_answer, the **correct** YES/NO uses a vivid blue (0,145,255); otherwise both are black.\n"
            if mark_answer
            else "Use --mark_answer to highlight the correct YES/NO in vivid **blue** (for review only).\n"
        )
    )
    _build_title(
        prs,
        f"Binary cue match — {' + '.join(rlist)}",
        sub,
    )
    for rk in rlist:
        items = list(load_binary_for_robot(rk))
        if sample_n is not None and len(items) > sample_n:
            rng = random.Random(seed)
            items = list(items)
            rng.shuffle(items)
            items = items[:sample_n]
        st = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(
            st, Inches(0.5), Inches(2.4), Inches(12.0), Inches(0.6),
            f"Section: {_labels(rk)}  ({len(items)} items)",
            size=20,
            bold=True,
        )
        total = len(items)
        appendix: list[list[str]] = []
        for o, it in enumerate(items, start=1):
            _slide_item(
                prs,
                o,
                total,
                it,
                _labels(rk),
                mark_answer=mark_answer,
                appendix=appendix,
                prefer_nr_gif=(rk == "manipulator"),
            )
        if items:
            n_before = len(prs.slides)
            n_app = max(1, (len(appendix) + 13) // 14)
            add_humaneval_appendix_slides(
                prs,
                title=f"Run index (binary_classification — {rk})",
                column_headers=APPENDIX_HEADERS,
                data_rows=appendix,
                total_slides_line=f"Total slides in this deck: {n_before + n_app}  (1-based; includes appendix).",
            )
    out_p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_p))
    print(str(out_p))
    return out_p


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--robot", type=str, default="all")
    ap.add_argument("--sample_n", type=str, default="None", help="max items per robot, or 'all'")
    ap.add_argument("--seed", type=int, default=17)
    ap.add_argument("--out", type=str, default=None)
    add_mark_answer_arg(ap)
    add_do_open_arg(ap)
    args = ap.parse_args()
    sn: int | None
    if args.sample_n in ("all", "None", "", "none"):
        sn = None
    else:
        sn = int(args.sample_n)
    out = build(
        robot=args.robot,
        sample_n=sn,
        seed=args.seed,
        out=args.out,
        mark_answer=mark_answer_effective_from_args(args),
    )
    open_output_file(out, do_open=do_open_effective_from_args(args))


if __name__ == "__main__":
    main()
