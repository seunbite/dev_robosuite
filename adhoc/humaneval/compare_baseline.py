#!/usr/bin/env python3
"""
Baseline comparison PPTX. **IIWA:** four **blinded** options A–D in one row
(Soph / No-reasoning / joint / xyzθ) — **multiselect** task: choose **all** motions
that match the cue. English prompt, cue in a **large rounded box** (like
binary_classification). When `--sample_n` is set, rows are drawn first from items
with all four baselines, then the rest, using `--position_seed` as the sampling
seed. A–D order uses `--position_seed` mixed per cue. `--mark_answer` colors A–D
only (**green** = sophisticated, **blue** = no-reasoning). GIF row is offset below
labels so A–D do not overlap the clips.

**TIAGo / quadruped:** two columns (unchanged).

Ends with a **Run index** appendix: deck slide #, item #, section, cue, index,
contextual, answer key. Total slide count in deck includes that appendix.

Output: data/pptx/YYYYMMDD_compare_baseline_<robot>.pptx
"""
from __future__ import annotations

import argparse
import hashlib
import random
import sys
from pathlib import Path

_H = Path(__file__).resolve().parent
if str(_H) not in sys.path:
    sys.path.insert(0, str(_H))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from _pptx_lib import (
    ANSWER_KEY_FONT_RGB,
    SOPH_KEY_GREEN_RGB,
    WS,
    MANIP_CFG,
    add_do_open_arg,
    add_humaneval_appendix_slides,
    add_mark_answer_arg,
    add_textbox,
    apply_sample_n,
    apply_sample_n_iiwa_fewshot_first,
    default_out_path,
    do_open_effective_from_args,
    fit_picture,
    latest_gif_in_dir,
    latest_iiwa_direct_baseline_gif,
    load_catalog_gifs,
    load_manipulator_product_list,
    mark_answer_effective_from_args,
    open_output_file,
    parse_robot,
)

_BLACK = RGBColor(0, 0, 0)


def _per_slide_shuffle_seed(cue: str, idx: int, sub: str) -> int:
    h = hashlib.md5(f"{sub}|{idx}|{cue}".encode()).hexdigest()
    return int(h[:8], 16) % (2**31)


def _iiwa_combined_shuffle_seed(position_seed: int, cue: str, idx: int, sub: str) -> int:
    """Tie global `--position_seed` to per-cue mix so shuffles are fixed but vary by slide."""
    base = int(position_seed) & 0x7FFFFFFF
    mix = _per_slide_shuffle_seed(cue, idx, sub)
    return (base * 0x9E3779B9 + mix) & 0x7FFFFFFF


IIWA_Q_EN = (
    "Select all options (A–D) whose motion, in your view, best captures the meaning of the given cue. "
    "You may select more than one. On paper or forms, list every letter that applies (e.g. A, C)."
)

APPENDIX_HEADERS = [
    "Deck slide #",
    "Item #",
    "Robot / section",
    "Cue name",
    "Cue index",
    "Contextual",
    "Answer key",
]


def _add_compare_slide(
    prs,
    title_line: str,
    desc: str,
    left_p: Path | None,
    right_p: Path | None,
    a_label: str,
    b_label: str,
    *,
    a_label_rgb: tuple[int, int, int] = (0, 0, 0),
    b_label_rgb: tuple[int, int, int] = (0, 0, 0),
) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, Inches(0.45), Inches(0.2), Inches(12.2), Inches(0.4), title_line, size=20, bold=True)
    if desc:
        add_textbox(s, Inches(0.5), Inches(0.6), Inches(12.0), Inches(0.55), desc[:500], size=10)
    gap = 0.15
    w = (13.333 - 2 * 0.45 - gap) / 2
    top = 1.2
    h = 5.0
    add_textbox(
        s, Inches(0.45), Inches(top - 0.2), Inches(6.0), Inches(0.2), a_label, size=12, bold=True, font_rgb=a_label_rgb
    )
    add_textbox(
        s, Inches(0.45 + w + gap), Inches(top - 0.2), Inches(6.0), Inches(0.2), b_label, size=12, bold=True, font_rgb=b_label_rgb
    )
    la = Inches(0.45)
    ra = Inches(0.45 + w + gap)
    if left_p and left_p.exists():
        fit_picture(s, left_p, la, Inches(top), Inches(w), Inches(h))
    else:
        add_textbox(s, la, Inches(top), Inches(w), Inches(0.35), f"(no GIF) {a_label}", size=9)
    if right_p and right_p.exists():
        fit_picture(s, right_p, ra, Inches(top), Inches(w), Inches(h))
    else:
        add_textbox(
            s, ra, Inches(top + 1.2), Inches(w), Inches(0.5),
            f"{b_label}\n(placeholder: add second render when available)",
            size=10,
        )


def _place_img_or_msg(slide, left, top, w, h, path: Path | None, missing: str) -> None:
    if path and path.exists():
        fit_picture(slide, path, left, top, w, h)
    else:
        add_textbox(slide, left, top, w, h, missing, size=7)


def _letter_color_mark(kind: str, mark_answer: bool) -> tuple[int, int, int]:
    if not mark_answer:
        return (0, 0, 0)
    if kind == "soph":
        return SOPH_KEY_GREEN_RGB
    if kind == "nr":
        return ANSWER_KEY_FONT_RGB
    return (0, 0, 0)


def _add_iiwa_baseline_row_slide(
    prs,
    subtest: str,
    idx: int,
    cue: str,
    desc: str,
    g_soph: Path | None,
    g_nr: Path | None,
    g_joint: Path | None,
    g_xyz: Path | None,
    mark_answer: bool,
    *,
    shuffle_seed: int,
    item_order: int,
) -> tuple[int, str]:
    """A–D in one row; letter + GIF only (policy names hidden). Returns (deck_slide_number, letter for sophisticated)."""
    letters = ("A", "B", "C", "D")
    panels: list[tuple[str, Path | None]] = [
        ("soph", g_soph),
        ("nr", g_nr),
        ("joint", g_joint),
        ("xyz", g_xyz),
    ]
    rng = random.Random(int(shuffle_seed))
    order = list(panels)
    rng.shuffle(order)
    soph_i = next(i for i, (k, _) in enumerate(order) if k == "soph")
    soph_letter = letters[soph_i]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    deck_slide = len(prs.slides)
    add_textbox(s, Inches(0.45), Inches(0.2), Inches(12.2), Inches(0.62), IIWA_Q_EN, size=13, bold=True)
    add_textbox(
        s, Inches(0.55), Inches(0.88), Inches(11.5), Inches(0.22),
        f"{subtest}  ·  c{idx:02d}  ·  item {item_order}",
        size=10,
        bold=True,
    )
    box = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65), Inches(1.12), Inches(11.9), Inches(0.55),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    box.line.color.rgb = RGBColor(200, 206, 216)
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = f"Cue: {cue}"
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = _BLACK
    if desc:
        add_textbox(s, Inches(0.5), Inches(1.72), Inches(12.0), Inches(0.45), desc[:600], size=8)
    m = 0.32
    gap = 0.07
    u = 13.333 - 2 * m
    cw = (u - 3 * gap) / 4
    lab_h = 0.26
    img_gap = 0.12
    ch = 2.52
    y0 = 2.18
    y_img = y0 + lab_h + img_gap
    for i in range(4):
        kind, g = order[i]
        letter = letters[i]
        col = _letter_color_mark(kind, mark_answer)
        x_off = m + i * (cw + gap)
        add_textbox(
            s, Inches(x_off), Inches(y0), Inches(cw), Inches(lab_h), letter, size=20, bold=True, font_rgb=col
        )
        _place_img_or_msg(
            s, Inches(x_off), Inches(y_img), Inches(cw), Inches(ch), g, f"{letter}: no GIF"
        )
    return deck_slide, soph_letter


def _manipulator_slides(
    prs,
    sample_n: int | None,
    mark_answer: bool,
    position_seed: int,
    appendix_rows: list[list[str]],
) -> None:
    st = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(
        st, Inches(0.5), Inches(2.0), Inches(12.0), Inches(0.95),
        "Section: Manipulator (IIWA) — one row A–D, blinded, **multiselect** (all valid options). "
        "Sampling: prefer items with all four renders (soph, nr, joint few-shot, xyz few-shot), then fill. "
        "Shuffle: `--position_seed` + per-cue hash (reproducible). "
        "mark_answer: **green** = sophisticated (include in any correct set), **blue** = no-reasoning.",
        size=11,
        bold=True,
    )
    rows = apply_sample_n_iiwa_fewshot_first(
        load_manipulator_product_list(), sample_n, position_seed
    )
    total = len(rows)
    for o, r in enumerate(rows, start=1):
        spec = MANIP_CFG[r.subtest]
        cue = r.cue
        idx = int(r.idx)
        g_soph = r.gif
        g_nr = latest_gif_in_dir(spec["nr_gif"], cue)
        g_joint = latest_iiwa_direct_baseline_gif(cue, idx, r.subtest, joint=True)
        g_xyz = latest_iiwa_direct_baseline_gif(cue, idx, r.subtest, joint=False)
        sh_seed = _iiwa_combined_shuffle_seed(position_seed, cue, idx, r.subtest)
        deck_sl, letter = _add_iiwa_baseline_row_slide(
            prs,
            r.subtest,
            idx,
            cue,
            r.label,
            g_soph,
            g_nr,
            g_joint,
            g_xyz,
            mark_answer=mark_answer,
            shuffle_seed=sh_seed,
            item_order=o,
        )
        ctx = "yes" if (r.subtest or "").lower() == "contextual" else "no"
        appendix_rows.append(
            [
                str(deck_sl),
                f"{o}/{total}",
                "IIWA",
                cue,
                str(idx),
                ctx,
                f"Must include: {letter} (soph.); multiselect — reviewers pick all that apply",
            ]
        )


def _embodiment_single_panel(
    prs,
    rk: str,
    label: str,
    sample_n: int | None,
    mark_answer: bool,
    appendix_rows: list[list[str]],
) -> None:
    st = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(
        st, Inches(0.5), Inches(2.0), Inches(12.0), Inches(0.8),
        f"Section: {label}\n(Only one render in repo; right column reserved for a future baseline.)",
        size=15,
        bold=True,
    )
    rows = apply_sample_n(load_catalog_gifs(rk), sample_n)
    total = len(rows)
    for o, r in enumerate(rows, start=1):
        _add_compare_slide(
            prs,
            f"{r.subtest} c{r.idx}  {r.cue}",
            r.label,
            r.gif,
            None,
            "Current (only) render",
            "Baseline (TBD)",
            a_label_rgb=ANSWER_KEY_FONT_RGB if mark_answer else (0, 0, 0),
            b_label_rgb=(0, 0, 0),
        )
        deck_sl = len(prs.slides)
        ctx = "yes" if (r.subtest or "").lower() == "contextual" else "no"
        ans = "Left column = current render (review key)" if mark_answer else "—"
        appendix_rows.append(
            [
                str(deck_sl),
                f"{o}/{total}",
                label.split("(")[0].strip(),
                r.cue,
                str(r.idx),
                ctx,
                ans,
            ]
        )


def build(
    robot: str = "all",
    sample_n: int | None = None,
    out: str | None = None,
    mark_answer: bool = False,
    position_seed: int = 20260424,
) -> Path:
    rlist = parse_robot(robot)
    out_p = (
        Path(out)
        if out
        else default_out_path("compare_baseline", "_".join(rlist) if len(rlist) > 1 else rlist[0])
    )
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(
        s0, Inches(0.5), Inches(0.5), Inches(12.0), Inches(0.5),
        "Best vs. baseline (side by side)",
        size=24,
        bold=True,
    )
    add_textbox(
        s0, Inches(0.6), Inches(1.1), Inches(12.0), Inches(1.6),
        f"IIWA: four options **A–D in one row**, **multiselect** (choose all that match the cue); policy names not shown. "
        f"Order is shuffled with RNG seed = `--position_seed` mixed per cue. "
        f"mark_answer: **green** = sophisticated, **blue** = no-reasoning. "
        f"Gold standard includes at least the sophisticated option; other picks are allowed.\n"
        f"TIAGo / quadruped: two columns; mark_answer: **Current** in blue.\n"
        f"{WS}",
        size=10,
    )
    appendix_rows: list[list[str]] = []
    for rk in rlist:
        if rk == "manipulator":
            _manipulator_slides(
                prs, sample_n, mark_answer=mark_answer, position_seed=position_seed, appendix_rows=appendix_rows
            )
        elif rk == "tiago":
            _embodiment_single_panel(
                prs, "tiago", "TIAGo (mobile_mani_v1)", sample_n, mark_answer, appendix_rows
            )
        else:
            _embodiment_single_panel(
                prs, "quadruped", "Quadruped (locomotion v1)", sample_n, mark_answer, appendix_rows
            )
    n_before_appendix = len(prs.slides)
    n_app = 1
    if appendix_rows:
        n_app = max(1, (len(appendix_rows) + 13) // 14)
    total_deck = n_before_appendix + n_app
    add_humaneval_appendix_slides(
        prs,
        title="Run index (compare_baseline)",
        column_headers=APPENDIX_HEADERS,
        data_rows=appendix_rows,
        total_slides_line=f"Total slides in this deck: {total_deck}  (1-based; includes this appendix).",
    )
    out_p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_p))
    print(str(out_p))
    return out_p


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--robot", type=str, default="all")
    ap.add_argument("--sample_n", type=str, default="None")
    ap.add_argument("--out", type=str, default=None)
    ap.add_argument(
        "--position_seed",
        type=int,
        default=20260424,
        help="Base RNG seed for A–D order (reproducible; mixed with per-cue hash per slide).",
    )
    add_mark_answer_arg(ap)
    add_do_open_arg(ap)
    a = ap.parse_args()
    sn: int | None
    if a.sample_n in ("all", "None", "", "none"):
        sn = None
    else:
        sn = int(a.sample_n)
    out = build(
        robot=a.robot,
        sample_n=sn,
        out=a.out,
        mark_answer=mark_answer_effective_from_args(a),
        position_seed=a.position_seed,
    )
    open_output_file(out, do_open=do_open_effective_from_args(a))


if __name__ == "__main__":
    main()
