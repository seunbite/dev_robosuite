#!/usr/bin/env python3
"""
all_product — Iconic → Contextual, full (or sampled) set of result GIFs in a grid per slide.

- per_page: max GIFs per slide (laid out in a rough sqrt grid)
- sample_n: max items per **robot** section, or None for all
- robot: manipulator | tiago | quadruped | all

Output: data/pptx/YYYYMMDD_all_product_<robot>.pptx
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

_H = Path(__file__).resolve().parent
if str(_H) not in sys.path:
    sys.path.insert(0, str(_H))

from pptx import Presentation
from pptx.util import Inches

from _pptx_lib import (
    WS,
    add_do_open_arg,
    add_humaneval_appendix_slides,
    add_textbox,
    default_out_path,
    do_open_effective_from_args,
    fit_picture,
    load_catalog_gifs,
    open_output_file,
    parse_robot,
    product_grid_layout,
    apply_sample_n,
)

APPENDIX_HEADERS = [
    "Deck slide #",
    "Item #",
    "Robot / section",
    "Cue name",
    "Cue index",
    "Contextual",
    "Answer key",
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = 0.45
HEADER = 0.48
GAP = 0.1


def _title_slide(prs, title: str, body: str) -> None:
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    add_textbox(s, Inches(M), Inches(0.35), Inches(12.0), Inches(0.8), title, size=24, bold=True)
    add_textbox(s, Inches(M), Inches(1.1), Inches(12.0), Inches(1.2), body, size=12)


def _section_title(prs, text: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, Inches(M), Inches(2.2), Inches(12.0), Inches(0.6), text, size=28, bold=True)


def _add_grid_slide(
    prs,
    chunk: list,
    *,
    per_page: int,
    robot_name: str,
    appendix: list[list[str]],
    item_index_base: int,
    total_items: int,
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    deck_sl = len(prs.slides)
    cols, _ = product_grid_layout(per_page)
    u_w = 13.333 - 2 * M
    u_h = 7.5 - M - 0.25 - HEADER
    cw = (u_w - (cols - 1) * GAP) / cols
    n_in = min(len(chunk), per_page)
    nrows = max(1, (n_in + cols - 1) // cols)
    ch = (u_h - (nrows - 1) * GAP) / nrows
    add_textbox(
        slide, Inches(M), Inches(0.12), Inches(12.0), Inches(0.3),
        f"{robot_name}  ·  iconic → contextual  ·  {len(chunk)} on this page",
        size=9,
    )
    for k, pr in enumerate(chunk):
        if k >= per_page:
            break
        row, col = divmod(k, cols)
        left = M + col * (cw + GAP)
        top = M + HEADER + row * (ch + GAP)
        cap_h = min(0.36, ch * 0.22)
        img_h = max(0.1, ch - cap_h)
        cline = f"{pr.subtest} c{pr.idx}  {pr.cue[:52]}"
        add_textbox(slide, Inches(left), Inches(top), Inches(cw), Inches(cap_h), cline, size=8, bold=True)
        if pr.gif and pr.gif.exists():
            fit_picture(slide, pr.gif, Inches(left), Inches(top + cap_h), Inches(cw), Inches(img_h))
        else:
            add_textbox(slide, Inches(left), Inches(top + cap_h), Inches(cw), Inches(0.28), f"(no gif) {pr.cue[:40]}", size=7)
        io = item_index_base + k + 1
        ctx = "yes" if (pr.subtest or "").lower() == "contextual" else "no"
        appendix.append(
            [
                str(deck_sl),
                f"{io}/{total_items}",
                robot_name,
                pr.cue[:200],
                str(pr.idx),
                ctx,
                "Grid all_product (order: iconic → contextual in catalog)",
            ]
        )


def build(
    robot: str = "all",
    per_page: int = 4,
    sample_n: int | None = None,
    out: str | None = None,
) -> Path:
    if per_page < 1:
        raise ValueError("per_page must be >= 1")
    rlist = parse_robot(robot)
    out_p = Path(out) if out else default_out_path("all_product", "_".join(rlist) if len(rlist) > 1 else rlist[0])
    out_p.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    rtag = " / ".join(rlist)
    _title_slide(
        prs,
        f"All product — {rtag}",
        f"per_page={per_page}  sample_n={sample_n!s}  (iconic first, then contextual within each section)\n"
        f"Workspace: {WS}",
    )

    for rk in rlist:
        if rk == "manipulator":
            label = "Manipulator (IIWA · no-reasoning v19)"
        elif rk == "tiago":
            label = "TIAGo (mobile_mani_v1)"
        else:
            label = "Quadruped (locomotion v1)"
        _section_title(prs, label)
        rows = apply_sample_n(
            load_catalog_gifs(rk, manipulator_gif="nr" if rk == "manipulator" else "soph"),
            sample_n,
        )
        appendix: list[list[str]] = []
        nt = len(rows)
        for i in range(0, len(rows), per_page):
            _add_grid_slide(
                prs,
                rows[i : i + per_page],
                per_page=per_page,
                robot_name=label,
                appendix=appendix,
                item_index_base=i,
                total_items=nt,
            )
        if appendix:
            n_before = len(prs.slides)
            n_app = max(1, (len(appendix) + 13) // 14)
            add_humaneval_appendix_slides(
                prs,
                title=f"Run index (all_product — {rk})",
                column_headers=APPENDIX_HEADERS,
                data_rows=appendix,
                total_slides_line=f"Total slides in this deck: {n_before + n_app}  (1-based; includes appendix).",
            )

    prs.save(str(out_p))
    print(str(out_p))
    return out_p


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--robot", type=str, default="all", help="manipulator, tiago, quadruped, or all")
    ap.add_argument("--per_page", type=int, default=4, help="max GIFs per slide (grid)")
    ap.add_argument(
        "--sample_n",
        type=str,
        default="None",
        help='max items per robot section (use "all" or None for full catalog)',
    )
    ap.add_argument("--out", type=str, default=None, help="output .pptx path")
    add_do_open_arg(ap)
    args = ap.parse_args()
    sn = None if args.sample_n in ("all", "None", "", "none", None) else int(args.sample_n)
    out = build(robot=args.robot, per_page=args.per_page, sample_n=sn, out=args.out)
    open_output_file(out, do_open=do_open_effective_from_args(args))


if __name__ == "__main__":
    main()
