#!/usr/bin/env python3
"""
One cue per slide: columns from --robot.
  - all: IIWA (soph) | TIAGo | quadruped, same cue name side by side
  - subset: that subset only

Order: manipulator (iconic → contextual) if it is in --robot, else tiago, else quadruped.

Output: data/pptx/YYYYMMDD_across_embodiment_<robot>.pptx
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

_H = Path(__file__).resolve().parent
if str(_H) not in sys.path:
    sys.path.insert(0, str(_H))

from pptx import Presentation
from pptx.util import Inches

from _pptx_lib import (
    WS,
    ProductRow,
    add_do_open_arg,
    add_humaneval_appendix_slides,
    add_textbox,
    apply_sample_n,
    default_out_path,
    do_open_effective_from_args,
    fit_picture,
    load_catalog_gifs,
    load_manipulator_product_list,
    open_output_file,
    parse_robot,
)

APPENDIX_HEADERS = [
    "Deck slide #",
    "Item #",
    "Robot / section",
    "Cue name",
    "Cue index",
    "Contextual",
    "Answer key",
]

LABELS = {
    "manipulator": "IIWA (no-reasoning v19)",
    "tiago": "TIAGo (mobile_mani_v1)",
    "quadruped": "Quadruped (locomotion v1)",
}


def _by_cue(rows: list[ProductRow]) -> dict[str, ProductRow]:
    return {r.cue: r for r in rows}


def _place_img(slide, path, left, top, w, h) -> None:
    if path and path.exists():
        fit_picture(slide, path, left, top, w, h)
    else:
        add_textbox(slide, left, top, w, Inches(0.35), "N/A (no GIF for this cue)", size=8)


def _ordered_rows(rlist: list[str], sample_n: int | None) -> list[ProductRow]:
    if "manipulator" in rlist:
        return apply_sample_n(load_manipulator_product_list(motion="nr"), sample_n)
    if "tiago" in rlist:
        return apply_sample_n(load_catalog_gifs("tiago"), sample_n)
    return apply_sample_n(load_catalog_gifs("quadruped"), sample_n)


def _gif_for(rk: str, cue: str, m_map: dict, t_map: dict, q_map: dict) -> Path | None:
    if rk == "manipulator":
        x = m_map.get(cue)
        return x.gif if x else None
    if rk == "tiago":
        x = t_map.get(cue)
        return x.gif if x else None
    x = q_map.get(cue)
    return x.gif if x else None


def build(robot: str = "all", sample_n: int | None = None, out: str | None = None) -> Path:
    rlist = parse_robot(robot)
    out_p = Path(out) if out else default_out_path("across_embodiment", "_".join(rlist) if len(rlist) > 1 else rlist[0])
    m_map = _by_cue(load_manipulator_product_list(motion="nr"))
    t_map = _by_cue(load_catalog_gifs("tiago"))
    q_map = _by_cue(load_catalog_gifs("quadruped"))
    m_rows = _ordered_rows(rlist, sample_n)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(
        s0, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.4),
        f"Same cue, multiple embodiments: {' | '.join(rlist)}",
        size=20,
        bold=True,
    )
    add_textbox(s0, Inches(0.55), Inches(0.95), Inches(12.0), Inches(0.3), f"{WS}", size=9)

    n = len(rlist)
    gap = 0.1
    m = 0.45
    u_w = 13.333 - 2 * m
    col_w = (u_w - (n - 1) * gap) / n
    top_lbl = 0.85
    top_img = 1.18
    h_img = 4.8

    appendix: list[list[str]] = []
    n_rows = len(m_rows)
    for o, r in enumerate(m_rows, start=1):
        cue = r.cue
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(
            s, Inches(m), Inches(0.2), Inches(12.0), Inches(0.4),
            f"{r.subtest}  ·  c{r.idx}  ·  {cue}",
            size=16,
            bold=True,
        )
        add_textbox(s, Inches(m), Inches(0.58), Inches(12.0), Inches(0.45), (r.label or "")[:500], size=9)
        for i, rk in enumerate(rlist):
            left = m + i * (col_w + gap)
            add_textbox(
                s, Inches(left), Inches(top_lbl - 0.1), Inches(col_w), Inches(0.22), LABELS[rk], size=9, bold=True
            )
            p = _gif_for(rk, cue, m_map, t_map, q_map)
            _place_img(s, p, Inches(left), Inches(top_img), Inches(col_w), Inches(h_img))
        ctx = "yes" if (r.subtest or "").lower() == "contextual" else "no"
        appendix.append(
            [
                str(len(prs.slides)),
                f"{o}/{n_rows}",
                " | ".join(rlist),
                cue,
                str(r.idx),
                ctx,
                "Same-cue across columns (order: manipulator, tiago, quadruped as selected)",
            ]
        )

    if appendix:
        n_before = len(prs.slides)
        n_app = max(1, (len(appendix) + 13) // 14)
        add_humaneval_appendix_slides(
            prs,
            title="Run index (across_embodiment)",
            column_headers=APPENDIX_HEADERS,
            data_rows=appendix,
            total_slides_line=f"Total slides in this deck: {n_before + n_app}  (1-based; includes appendix).",
        )

    out_p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_p))
    print(str(out_p))
    return out_p


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--robot", type=str, default="all")
    ap.add_argument("--sample_n", type=str, default="None")
    ap.add_argument("--out", type=str, default=None)
    add_do_open_arg(ap)
    a = ap.parse_args()
    sn: int | None
    if a.sample_n in ("all", "None", "", "none"):
        sn = None
    else:
        sn = int(a.sample_n)
    out = build(robot=a.robot, sample_n=sn, out=a.out)
    open_output_file(out, do_open=do_open_effective_from_args(a))


if __name__ == "__main__":
    main()
