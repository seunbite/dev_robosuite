"""
Compact PPTX summary for binary_classification VLM JSONL (one slide per instance × robot).
"""
from __future__ import annotations

import json
from collections import defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

# Order must match expand_input_types ``all``
COL_ORDER: list[tuple[str, str]] = [
    ("mp4", "mp4"),
    ("alpha_frame", "alpha"),
    ("first_frame_trajectory", "first+traj"),
    ("alpha_frame_trajectory", "alpha+traj"),
    ("mp4_plus_trajectory", "mp4+traj"),
]


def build_binary_results_pptx_from_records(records: list[dict[str, Any]], out_pptx: Path) -> Path:
    grouped: dict[tuple[str, str], dict[str, str]] = defaultdict(dict)
    order: list[tuple[str, str]] = []
    for r in records:
        if r.get("task") != "binary_classification":
            continue
        rob = str(r.get("robot", ""))
        iid = str(r.get("instance_id", ""))
        it = str(r.get("input_type", ""))
        key = (rob, iid)
        if key not in order:
            order.append(key)
        grouped[key][it] = (r.get("response") or "").strip()

    prs = Presentation()
    layout = prs.slide_layouts[5]

    for rob, iid in order:
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        title.text = f"{rob} — {iid}"
        if title.text_frame.paragraphs:
            title.text_frame.paragraphs[0].font.size = Pt(13)

        rows, cols = 2, len(COL_ORDER)
        left = Inches(0.2)
        top = Inches(0.95)
        width = Inches(9.5)
        height = Inches(6.0)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for j, (it, short) in enumerate(COL_ORDER):
            c0 = table.cell(0, j)
            c0.text = short
            for p in c0.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
            txt = grouped[(rob, iid)].get(it, "—")
            c1 = table.cell(1, j)
            c1.text = txt[:6000]
            for p in c1.text_frame.paragraphs:
                p.font.size = Pt(6.5)

    out_pptx = Path(out_pptx)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    return out_pptx.resolve()


def build_binary_results_pptx_from_jsonl(jsonl: Path, out_pptx: Path) -> Path:
    lines = Path(jsonl).read_text(encoding="utf-8").splitlines()
    recs = [json.loads(l) for l in lines if l.strip()]
    return build_binary_results_pptx_from_records(recs, out_pptx)
