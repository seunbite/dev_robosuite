from __future__ import annotations

import json
import random
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image, ImageSequence


ROOT = Path("/Users/sb/Downloads/workspace/dev_robosuite")
SEED_ROOT = ROOT / "data" / "seed"
MOTION_ROOT = ROOT / "data" / "motions"
OUT_ROOT = SEED_ROOT / "goodset_human_labeling"
LITE_ASSET_ROOT = OUT_ROOT / "lite_assets_320"

ICONIC_PPT = SEED_ROOT / "iconic_nr.pptx"
CONTEXTUAL_PPT = SEED_ROOT / "contextual_nr.pptx"

SOPH_ROOTS = {
    "iconic": MOTION_ROOT / "v19_sophisticated" / "IIWA",
    "contextual": MOTION_ROOT / "v19_sophisticated_contextual" / "IIWA",
}
NR_ROOTS = {
    "iconic": MOTION_ROOT / "baseline_prompt19_full_no_reasoning" / "no_reasoning_iconic" / "IIWA",
    "contextual": MOTION_ROOT / "baseline_prompt19_full_no_reasoning" / "no_reasoning_contextual" / "IIWA",
}
DIRECT_ROOTS = [
    MOTION_ROOT / "baseline_prompt19_direct_experiment",
    MOTION_ROOT / "baseline_prompt19_direct_experiment_extra10",
    MOTION_ROOT / "baseline_prompt19_direct_experiment_nonoverlap10",
    MOTION_ROOT / "baseline_prompt19_direct_experiment_more20configs",
]

FILENAME_RE = re.compile(r"(?P<ts>\d{8}_\d{6})_IIWA_(?P<cue>.+?)_p\d+\.gif$")
TYPE_COLORS = {
    "pose": RGBColor(34, 197, 94),
    "movement": RGBColor(59, 130, 246),
    "path": RGBColor(139, 92, 246),
    "gripper": RGBColor(245, 158, 11),
    "joint": RGBColor(239, 68, 68),
    "xyz": RGBColor(20, 184, 166),
}


@dataclass
class CueItem:
    dataset: str
    idx: int
    cue: str
    description: str


def load_ppt_good_bad(path: Path, dataset: str) -> Dict[str, List[CueItem]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    mode: Optional[str] = None
    items = {"good": [], "bad": []}
    for slide in prs.slides:
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        joined = "\n".join(texts)
        if "Good Examples" in joined:
            mode = "good"
            continue
        if "Failures" in joined:
            mode = "bad"
            continue
        m = re.search(r"c(\d+)\s+([a-z0-9_]+)", joined)
        if not m or not mode:
            continue
        idx = int(m.group(1))
        cue = m.group(2)
        description = ""
        if len(texts) >= 2:
            description = texts[1].replace("\n", " ").strip()
        items[mode].append(CueItem(dataset=dataset, idx=idx, cue=cue, description=description))
    return items


def token_set(cue: str) -> set[str]:
    return set(part for part in cue.split("_") if part)


def choose_negative(cue: CueItem, pool: Iterable[CueItem], rng: random.Random) -> CueItem:
    src_tokens = token_set(cue.cue)
    candidates = [
        item
        for item in pool
        if item.cue != cue.cue and not (src_tokens & token_set(item.cue))
    ]
    if not candidates:
        candidates = [item for item in pool if item.cue != cue.cue]
    return rng.choice(sorted(candidates, key=lambda x: (x.idx, x.cue)))


def latest_gif_by_cue(root: Path) -> Dict[str, Path]:
    found: Dict[str, tuple[str, Path]] = {}
    if not root.exists():
        return {}
    for path in root.glob("*.gif"):
        m = FILENAME_RE.match(path.name)
        if not m:
            continue
        cue = m.group("cue")
        ts = m.group("ts")
        prev = found.get(cue)
        if prev is None or ts > prev[0]:
            found[cue] = (ts, path)
    return {cue: path for cue, (_, path) in found.items()}


def latest_direct_gifs() -> Dict[str, Dict[str, Dict[str, Path]]]:
    result: Dict[str, Dict[str, Dict[str, tuple[str, Path]]]] = {
        "joint": {"iconic": {}, "contextual": {}},
        "xyz": {"iconic": {}, "contextual": {}},
    }
    for root in DIRECT_ROOTS:
        for path in root.rglob("*.gif"):
            parts = path.parts
            kind = "joint" if "direct_joint" in parts else "xyz" if "direct_xyz_theta" in parts else None
            if not kind:
                continue
            dataset = "iconic" if "iconic" in parts else "contextual" if "contextual" in parts else None
            if not dataset:
                continue
            m = FILENAME_RE.match(path.name)
            if not m:
                continue
            cue = m.group("cue")
            ts = m.group("ts")
            prev = result[kind][dataset].get(cue)
            if prev is None or ts > prev[0]:
                result[kind][dataset][cue] = (ts, path)
    return {
        kind: {
            dataset: {cue: path for cue, (_, path) in per_ds.items()}
            for dataset, per_ds in per_kind.items()
        }
        for kind, per_kind in result.items()
    }


def _load_json_rows(path: Path) -> list[dict]:
    if not path.exists():
        return []
    return json.loads(path.read_text())


def load_config_maps() -> dict[str, dict[str, dict]]:
    maps = {
        "sophisticated": {"iconic": {}, "contextual": {}},
        "no_reasoning": {"iconic": {}, "contextual": {}},
    }
    for row in _load_json_rows(SEED_ROOT / "motion_configs_prompt_v19_sophisticated.json"):
        maps["sophisticated"]["iconic"][row["cue"]] = row
    for row in _load_json_rows(SEED_ROOT / "motion_configs_prompt_v19_sophisticated_contextual.json"):
        maps["sophisticated"]["contextual"][row["cue"]] = row
    for row in _load_json_rows(SEED_ROOT / "baseline_prompt19_full_no_reasoning" / "motion_configs_prompt_v19_sophisticated_no_reasoning_iconic.json"):
        maps["no_reasoning"]["iconic"][row["cue"]] = row
    for row in _load_json_rows(SEED_ROOT / "baseline_prompt19_full_no_reasoning" / "motion_configs_prompt_v19_sophisticated_no_reasoning_contextual.json"):
        maps["no_reasoning"]["contextual"][row["cue"]] = row
    return maps


def compress_gif_for_ppt(src: Path, group: str) -> Path:
    out_dir = LITE_ASSET_ROOT / group
    out_dir.mkdir(parents=True, exist_ok=True)
    dst = out_dir / src.name
    if dst.exists():
        return dst
    img = Image.open(src)
    frames = []
    durations = []
    loop = img.info.get("loop", 0)
    for frame in ImageSequence.Iterator(img):
        rgba = frame.convert("RGBA")
        scale = min(1.0, 320 / rgba.width)
        size = (max(1, int(rgba.width * scale)), max(1, int(rgba.height * scale)))
        resized = rgba.resize(size, Image.Resampling.LANCZOS)
        white = Image.new("RGBA", size, (255, 255, 255, 255))
        white.alpha_composite(resized)
        frames.append(white.convert("P", palette=Image.Palette.ADAPTIVE))
        durations.append(frame.info.get("duration", img.info.get("duration", 120)))
    frames[0].save(dst, save_all=True, append_images=frames[1:], duration=durations, loop=loop, optimize=True, disposal=2)
    return dst


def add_textbox(slide, left, top, width, height, text, size=20, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0, 0, 0)
    return box


def add_choice_chip(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(245, 247, 250)
    shape.line.color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(180, 186, 196)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = RGBColor(0, 0, 0)
    return shape


def fit_picture(slide, img_path: Path, left, top, width, height):
    from PIL import Image

    with Image.open(img_path) as im:
        img_w, img_h = im.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        draw_w = width
        draw_h = width / img_ratio
        draw_left = left
        draw_top = top + (height - draw_h) / 2
    else:
        draw_h = height
        draw_w = height * img_ratio
        draw_top = top
        draw_left = left + (width - draw_w) / 2
    slide.shapes.add_picture(str(img_path), draw_left, draw_top, draw_w, draw_h)


def build_binary_ppt(dataset: str, goods: List[CueItem], cue_pool: List[CueItem], nr_gifs: Dict[str, Path], out_path: Path, rng_seed: int, compress: bool = False) -> List[dict]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(title, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6), f"Prompt 19 {dataset.title()} Good Set: Binary Cue Matching", 24, True)
    add_textbox(
        title,
        Inches(0.8),
        Inches(1.3),
        Inches(11.8),
        Inches(2.0),
        "Task: Look at the IIWA motion and answer whether the shown cue matches the motion.\nChoose Yes if it matches, No if it does not.\nEach good cue appears twice overall: once with the true cue and once with a non-overlapping random false cue.",
        18,
    )
    add_choice_chip(title, Inches(1.0), Inches(3.2), Inches(1.5), Inches(0.55), "YES")
    add_choice_chip(title, Inches(2.8), Inches(3.2), Inches(1.5), Inches(0.55), "NO")

    rng = random.Random(rng_seed)
    manifest = []
    items = []
    ordered_goods = sorted(goods, key=lambda x: (x.idx, x.cue))
    for cue in ordered_goods:
        neg = choose_negative(cue, cue_pool, rng)
        items.append({
            "cue_item": cue,
            "shown_cue": cue.cue,
            "shown_description": cue.description,
            "ground_truth": "yes",
        })
        items.append({
            "cue_item": cue,
            "shown_cue": neg.cue,
            "shown_description": neg.description,
            "ground_truth": "no",
        })

    for order, item in enumerate(items, start=1):
        cue = item["cue_item"]
        gif = nr_gifs.get(cue.cue)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(slide, Inches(0.45), Inches(0.3), Inches(12.2), Inches(0.45), f"{dataset.title()} c{cue.idx}  {cue.cue}", 24, True)
        add_textbox(slide, Inches(0.6), Inches(0.85), Inches(12.0), Inches(0.5), "Question: Does the shown cue match this motion?", 17, True)
        cue_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(11.9), Inches(1.0))
        cue_box.fill.solid()
        cue_box.fill.fore_color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(248, 250, 252)
        cue_box.line.color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(200, 206, 216)
        tf = cue_box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = f"Shown cue: {item['shown_cue']}"
        r1.font.size = Pt(22)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(0, 0, 0)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = item["shown_description"]
        r2.font.size = Pt(15)
        r2.font.name = "Arial"
        r2.font.color.rgb = RGBColor(0, 0, 0)
        if gif and gif.exists():
            used_gif = compress_gif_for_ppt(gif, f"binary_{dataset}") if compress else gif
            fit_picture(slide, used_gif, Inches(2.0), Inches(2.55), Inches(9.3), Inches(3.75))
        add_choice_chip(slide, Inches(3.6), Inches(6.55), Inches(2.1), Inches(0.55), "YES")
        add_choice_chip(slide, Inches(7.1), Inches(6.55), Inches(2.1), Inches(0.55), "NO")
        add_textbox(slide, Inches(11.5), Inches(6.75), Inches(1.2), Inches(0.25), f"{order}/{len(items)}", 10, False, PP_ALIGN.RIGHT)
        manifest.append({
            "order": order,
            "dataset": dataset,
            "source_cue_idx": cue.idx,
            "source_cue": cue.cue,
            "shown_cue": item["shown_cue"],
            "ground_truth": item["ground_truth"],
            "gif": str((compress_gif_for_ppt(gif, f'binary_{dataset}') if compress and gif else gif)) if gif else None,
        })

    prs.save(str(out_path))
    return manifest


def build_fourway_ppt(dataset: str, goods: List[CueItem], assets: Dict[str, Dict[str, Path]], out_path: Path, rng_seed: int, compress: bool = False) -> List[dict]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(title, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.6), f"Prompt 19 {dataset.title()} Good Set: Pick The Best Motion", 24, True)
    add_textbox(
        title,
        Inches(0.8),
        Inches(1.3),
        Inches(11.8),
        Inches(2.0),
        "Task: For each cue, compare four IIWA motions and select all options that appropriately match the cue.\nThere may be more than one acceptable answer.\nOptions come from Sophisticated, No Reasoning, Direct Joint, and Direct XYZ-Theta when all four are available.",
        18,
    )
    add_choice_chip(title, Inches(1.0), Inches(3.2), Inches(1.2), Inches(0.55), "A")
    add_choice_chip(title, Inches(2.5), Inches(3.2), Inches(1.2), Inches(0.55), "B")
    add_choice_chip(title, Inches(4.0), Inches(3.2), Inches(1.2), Inches(0.55), "C")
    add_choice_chip(title, Inches(5.5), Inches(3.2), Inches(1.2), Inches(0.55), "D")

    rng = random.Random(rng_seed)
    manifest = []
    candidates = sorted([
        cue for cue in goods
        if cue.cue in assets["sophisticated"]
        and cue.cue in assets["no_reasoning"]
        and cue.cue in assets["joint"]
        and cue.cue in assets["xyz"]
    ], key=lambda x: (x.idx, x.cue))[:30]

    for order, cue in enumerate(candidates, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(slide, Inches(0.45), Inches(0.22), Inches(12.2), Inches(0.42), f"{dataset.title()} c{cue.idx}  {cue.cue}", 22, True)
        add_textbox(slide, Inches(0.55), Inches(0.65), Inches(12.0), Inches(0.65), cue.description, 14, False)
        add_textbox(slide, Inches(0.55), Inches(1.18), Inches(12.0), Inches(0.3), "Question: Select all options that appropriately match the cue.", 15, True)

        variants = [
            ("sophisticated", assets["sophisticated"][cue.cue]),
            ("no_reasoning", assets["no_reasoning"][cue.cue]),
            ("direct_joint", assets["joint"][cue.cue]),
            ("direct_xyz_theta", assets["xyz"][cue.cue]),
        ]
        rng.shuffle(variants)
        positions = [
            ("A", Inches(0.25), Inches(2.0)),
            ("B", Inches(3.45), Inches(2.0)),
            ("C", Inches(6.65), Inches(2.0)),
            ("D", Inches(9.85), Inches(2.0)),
        ]
        slide_record = {
            "order": order,
            "dataset": dataset,
            "cue_idx": cue.idx,
            "cue": cue.cue,
            "options": {},
        }
        for (label, left, top), (variant_name, gif_path) in zip(positions, variants):
            add_choice_chip(slide, left, top - Inches(0.42), Inches(0.6), Inches(0.34), label)
            used_gif = compress_gif_for_ppt(gif_path, f"{dataset}_{variant_name}") if compress else gif_path
            fit_picture(slide, used_gif, left, top, Inches(2.95), Inches(4.8))
            slide_record["options"][label] = {
                "variant": variant_name,
                "gif": str(used_gif),
            }
        add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.2), Inches(0.2), f"{order}/{len(candidates)}", 10, False, PP_ALIGN.RIGHT)
        manifest.append(slide_record)

    prs.save(str(out_path))
    return manifest


def main():
    OUT_ROOT.mkdir(parents=True, exist_ok=True)
    LITE_ASSET_ROOT.mkdir(parents=True, exist_ok=True)

    iconic = load_ppt_good_bad(ICONIC_PPT, "iconic")
    contextual = load_ppt_good_bad(CONTEXTUAL_PPT, "contextual")
    cue_pool = {
        "iconic": iconic["good"] + iconic["bad"],
        "contextual": contextual["good"] + contextual["bad"],
    }

    soph = {ds: latest_gif_by_cue(root) for ds, root in SOPH_ROOTS.items()}
    nr = {ds: latest_gif_by_cue(root) for ds, root in NR_ROOTS.items()}
    direct = latest_direct_gifs()
    manifest = {
        "good_sets": {
            "iconic": [item.__dict__ for item in iconic["good"]],
            "contextual": [item.__dict__ for item in contextual["good"]],
        },
        "bad_sets": {
            "iconic": [item.__dict__ for item in iconic["bad"]],
            "contextual": [item.__dict__ for item in contextual["bad"]],
        },
    }

    binary_iconic_path = OUT_ROOT / "prompt19_goodset_binary_iconic_iiwa_20260411_v3_sorted.pptx"
    binary_contextual_path = OUT_ROOT / "prompt19_goodset_binary_contextual_iiwa_20260411_v3_sorted.pptx"
    four_iconic_path = OUT_ROOT / "prompt19_goodset_pick_best_iconic_iiwa_20260411_v4_top30_1row.pptx"
    four_contextual_path = OUT_ROOT / "prompt19_goodset_pick_best_contextual_iiwa_20260411_v4_top30_1row.pptx"

    manifest["binary_iconic"] = build_binary_ppt(
        "iconic", iconic["good"], cue_pool["iconic"], nr["iconic"], binary_iconic_path, rng_seed=11
    )
    manifest["binary_contextual"] = build_binary_ppt(
        "contextual", contextual["good"], cue_pool["contextual"], nr["contextual"], binary_contextual_path, rng_seed=17
    )
    manifest["pick_best_iconic"] = build_fourway_ppt(
        "iconic",
        iconic["good"],
        {
            "sophisticated": soph["iconic"],
            "no_reasoning": nr["iconic"],
            "joint": direct["joint"]["iconic"],
            "xyz": direct["xyz"]["iconic"],
        },
        four_iconic_path,
        rng_seed=23,
        compress=False,
    )
    manifest["pick_best_contextual"] = build_fourway_ppt(
        "contextual",
        contextual["good"],
        {
            "sophisticated": soph["contextual"],
            "no_reasoning": nr["contextual"],
            "joint": direct["joint"]["contextual"],
            "xyz": direct["xyz"]["contextual"],
        },
        four_contextual_path,
        rng_seed=29,
        compress=False,
    )

    lite_four_iconic_path = OUT_ROOT / "prompt19_goodset_pick_best_iconic_iiwa_20260412_lite_320_bar.pptx"
    lite_four_contextual_path = OUT_ROOT / "prompt19_goodset_pick_best_contextual_iiwa_20260412_lite_320_bar.pptx"
    manifest["pick_best_iconic_lite"] = build_fourway_ppt(
        "iconic",
        iconic["good"],
        {
            "sophisticated": soph["iconic"],
            "no_reasoning": nr["iconic"],
            "joint": direct["joint"]["iconic"],
            "xyz": direct["xyz"]["iconic"],
        },
        lite_four_iconic_path,
        rng_seed=23,
        compress=True,
    )
    manifest["pick_best_contextual_lite"] = build_fourway_ppt(
        "contextual",
        contextual["good"],
        {
            "sophisticated": soph["contextual"],
            "no_reasoning": nr["contextual"],
            "joint": direct["joint"]["contextual"],
            "xyz": direct["xyz"]["contextual"],
        },
        lite_four_contextual_path,
        rng_seed=29,
        compress=True,
    )

    (OUT_ROOT / "manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print("WROTE", OUT_ROOT / "manifest.json")
    print("WROTE", binary_iconic_path)
    print("WROTE", binary_contextual_path)
    print("WROTE", four_iconic_path)
    print("WROTE", four_contextual_path)
    print("WROTE", lite_four_iconic_path)
    print("WROTE", lite_four_contextual_path)


if __name__ == "__main__":
    main()
