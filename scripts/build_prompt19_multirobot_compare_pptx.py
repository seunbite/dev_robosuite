from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path

import fire
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/sb/Downloads/workspace/dev_robosuite")
SEED = ROOT / "data" / "seed"
MOTIONS = ROOT / "data" / "motions"
OUT_DIR = SEED

ROBOTS = ["IIWA", "Panda", "XArm7"]

CONFIGS = {
    "iconic": {
        "sophisticated": SEED / "motion_configs_prompt_v19_sophisticated.json",
        "no_reasoning": SEED / "baseline_prompt19_full_no_reasoning" / "motion_configs_prompt_v19_sophisticated_no_reasoning_iconic.json",
        "sophisticated_gif_dir": lambda robot: MOTIONS / "v19_sophisticated" / robot,
        "no_reasoning_gif_dir": lambda robot: MOTIONS / "baseline_prompt19_full_no_reasoning" / "no_reasoning_iconic" / robot,
    },
    "contextual": {
        "sophisticated": SEED / "motion_configs_prompt_v19_sophisticated_contextual.json",
        "no_reasoning": SEED / "baseline_prompt19_full_no_reasoning" / "motion_configs_prompt_v19_sophisticated_no_reasoning_contextual.json",
        "sophisticated_gif_dir": lambda robot: (
            MOTIONS / ("v19_sophisticated_contextual_q4filled" if robot == "IIWA" else "v19_sophisticated_contextual") / robot
        ),
        "no_reasoning_gif_dir": lambda robot: MOTIONS / "baseline_prompt19_full_no_reasoning" / "no_reasoning_contextual" / robot,
    },
}


def _load_json(path: Path) -> list[dict]:
    return json.loads(path.read_text(encoding="utf-8"))


def _safe_name(text: str) -> str:
    return str(text).replace("/", "_").replace("\\", "_").replace(" ", "_")


def _latest_single_gif(base: Path, cue: str) -> Path | None:
    safe = _safe_name(cue)
    matches = sorted(base.rglob(f"*_{safe}_p*.gif"), key=lambda p: p.stat().st_mtime)
    return matches[-1] if matches else None


def _add_textbox(slide, left, top, width, height, text, *, font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return box


def _fit_desc(text: str, max_len: int = 220) -> str:
    text = (text or "").strip()
    if len(text) <= max_len:
        return text
    return text[: max_len - 1].rstrip() + "…"


def _add_image_or_placeholder(slide, left, top, width, height, path: Path | None, label: str):
    if path is not None and path.exists():
        with Image.open(path) as img:
            img_w, img_h = img.size
        if img_w > 0 and img_h > 0:
            box_w = int(width)
            box_h = int(height)
            scale = min(box_w / img_w, box_h / img_h)
            fitted_w = int(img_w * scale)
            fitted_h = int(img_h * scale)
            offset_x = int((box_w - fitted_w) / 2)
            offset_y = int((box_h - fitted_h) / 2)
            slide.shapes.add_picture(
                str(path),
                left + offset_x,
                top + offset_y,
                width=fitted_w,
                height=fitted_h,
            )
        else:
            slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        shape = slide.shapes.add_shape(1, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = None  # type: ignore[assignment]
        shape.text_frame.text = "No image"
    _add_textbox(slide, left, top - Inches(0.22), width, Inches(0.18), label, font_size=11, bold=True)


def build(
    dataset: str = "contextual",
    include_no_reasoning: bool = False,
    show_section_labels: bool = False,
    source_variant: str = "sophisticated",
    start_idx: int | None = None,
    end_idx: int | None = None,
    output_name: str | None = None,
) -> str:
    if dataset not in CONFIGS:
        raise ValueError(f"dataset must be one of {sorted(CONFIGS)}")
    if source_variant not in {"sophisticated", "no_reasoning"}:
        raise ValueError("source_variant must be one of: sophisticated, no_reasoning")

    spec = CONFIGS[dataset]
    selected_rows = sorted(_load_json(spec[source_variant]), key=lambda r: int(r["idx"]))
    if start_idx is not None:
        selected_rows = [r for r in selected_rows if int(r["idx"]) >= int(start_idx)]
    if end_idx is not None:
        selected_rows = [r for r in selected_rows if int(r["idx"]) <= int(end_idx)]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = f"Prompt 19 {dataset.title()} {'No Reasoning' if source_variant == 'no_reasoning' else 'Sophisticated'}"
    _add_textbox(title_slide, Inches(0.6), Inches(0.55), Inches(12.1), Inches(0.6), title, font_size=26, bold=True)
    subtitle = "IIWA / Panda / XArm7"
    if include_no_reasoning:
        subtitle = "Top row: IIWA / Panda / XArm7 | Bottom row: IIWA / Panda / XArm7"
    _add_textbox(title_slide, Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.6), subtitle, font_size=15)
    _add_textbox(
        title_slide,
        Inches(0.6),
        Inches(1.7),
        Inches(12.0),
        Inches(0.5),
        f"Generated {datetime.now().isoformat(timespec='seconds')}",
        font_size=12,
    )

    left_margin = Inches(0.42)
    gap = Inches(0.12)
    img_w = Inches(4.08)
    top_row_y = Inches(1.65)
    if include_no_reasoning:
        img_h = Inches(1.95)
        bottom_row_y = Inches(4.55)
    else:
        img_h = Inches(4.65)
        bottom_row_y = None

    for row in selected_rows:
        idx = int(row["idx"])
        cue = row["cue"]
        desc = _fit_desc(row.get("description", ""))

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, Inches(0.45), Inches(0.22), Inches(12.2), Inches(0.38), f"c{idx} {cue}", font_size=22, bold=True)
        _add_textbox(slide, Inches(0.45), Inches(0.68), Inches(12.2), Inches(0.55), desc, font_size=12)

        if include_no_reasoning and show_section_labels:
            _add_textbox(slide, Inches(0.45), Inches(1.36), Inches(6.0), Inches(0.22), "Sophisticated", font_size=14, bold=True)
            _add_textbox(slide, Inches(0.45), Inches(3.86), Inches(6.0), Inches(0.22), "No Reasoning", font_size=14, bold=True)

        for i, robot in enumerate(ROBOTS):
            left = left_margin + i * (img_w + gap)
            primary_gif_dir = spec["sophisticated_gif_dir"] if source_variant == "sophisticated" else spec["no_reasoning_gif_dir"]
            primary_gif = _latest_single_gif(primary_gif_dir(robot), cue)
            _add_image_or_placeholder(slide, left, top_row_y, img_w, img_h, primary_gif, robot)
            if include_no_reasoning and bottom_row_y is not None:
                nr_gif = _latest_single_gif(spec["no_reasoning_gif_dir"](robot), cue)
                _add_image_or_placeholder(slide, left, bottom_row_y, img_w, img_h, nr_gif, robot)

    if output_name is None:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        if include_no_reasoning:
            suffix = "multirobot_compare"
        else:
            suffix = f"multirobot_{source_variant}_only"
        range_suffix = ""
        if start_idx is not None or end_idx is not None:
            range_suffix = f"_c{start_idx if start_idx is not None else 'start'}_to_c{end_idx if end_idx is not None else 'end'}"
        output_name = f"prompt19_{dataset}_{suffix}{range_suffix}_{stamp}.pptx"

    out_path = OUT_DIR / output_name
    prs.save(str(out_path))
    print(out_path)
    return str(out_path)


if __name__ == "__main__":
    fire.Fire({"build": build})
